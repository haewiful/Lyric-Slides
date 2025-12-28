import tkinter as tk
from tkinter import messagebox, filedialog, font, ttk
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

class Lyrics(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("가사 파워포인트 생성 프로그램")

        self.ppt = Presentation()

        self.entry_widgets = {}
        # Contains entries named: [lyrics, english lyrics, font size, font family]

        self.create_frame()

        self.after(0, self.center_window)
    
    def create_frame(self):
        # Lyrics Frame
        self.lyrics_frame = tk.Frame(self)

        input_row_n=0
        columnspan_n=4
        self.lyrics_frame.grid(row=input_row_n, column=0, columnspan=3, padx=10, pady=10, sticky="nsew")
        self.lyrics_frame.grid_columnconfigure(1, weight=1)
        self.lyrics_frame.grid_columnconfigure(2, weight=1)

        input_row_n+=1

        tk.Label(self.lyrics_frame, text="가사 입력").grid(row=input_row_n, column=0, columnspan=columnspan_n)
        input_row_n+=1

        lyrics_text = tk.Text(self.lyrics_frame, width=50, height=20)
        lyrics_text.grid(row=input_row_n, column=0, columnspan=5, pady=5)
        self.entry_widgets["lyrics"] = lyrics_text
        input_row_n+=1

        ## font size entry
        font_default = 48
        tk.Label(self.lyrics_frame, text="Font Size:", anchor="e").grid(row=input_row_n, column=0, padx=0, pady=5, sticky="nsew")
        font_entry = tk.Entry(self.lyrics_frame)
        font_entry.insert(0, str(font_default))
        font_entry.grid(row=input_row_n, column=1, columnspan=columnspan_n-1, padx=0, pady=5, sticky="nsew")

        self.entry_widgets["font size"] = font_entry
        input_row_n+=1

        ## font family entry: Dynamically get the list of fonts from the system
        font_families = sorted(list(font.families()))

        tk.Label(self.lyrics_frame, text="Font:", anchor="e").grid(row=input_row_n, column=0, padx=0, pady=5, sticky="nsew")
        font_family_combobox = ttk.Combobox(self.lyrics_frame, values=font_families, state="readonly")
        font_family_combobox.grid(row=input_row_n, column=1, columnspan=columnspan_n-1, padx=0, pady=5, sticky="nsew")
        # Check if a common font is in the list to set as a default
        if "Arial" in font_families:
            font_family_combobox.set("Arial")
        else:
            font_family_combobox.set(font_families[0])

        self.entry_widgets["font family"] = font_family_combobox
        input_row_n+=1

        ## Number of lines per slide
        tk.Label(self.lyrics_frame, text="# of lines per slide:", anchor="e").grid(row=input_row_n, column=0, padx=0, pady=5, sticky="nsew")
        line_num_combobox = ttk.Combobox(self.lyrics_frame, textvariable="2", values=[str(y) for y in range(1, 5)], state="readonly")
        line_num_combobox.set(2)
        line_num_combobox.grid(row=input_row_n, column=1, columnspan=columnspan_n-1, padx=0, pady=5, sticky="nsew")
        self.entry_widgets["num of lines"] = line_num_combobox
        input_row_n+=1

        # lyrics position radio buttons
        tk.Label(self.lyrics_frame, text="가사 위치:", anchor="e").grid(row=input_row_n, column=0, padx=0, pady=5, sticky="nsew")
        position_options = ["Top", "Center", "Bottom"]
        self.position_var = tk.StringVar(value=position_options[1])
        for i, option in enumerate(position_options):
            rb = tk.Radiobutton(self.lyrics_frame, text=option, anchor="w", variable=self.position_var, value=option)
            rb.grid(row=input_row_n, column=i+1, padx=0, pady=5, sticky="nsew")
        self.entry_widgets["lyrics position"] = self.position_var
        input_row_n+=1

        ## Line position
        tk.Label(self.lyrics_frame, text="English lyrics?", anchor="e").grid(row=input_row_n, column=0, padx=0, pady=5, sticky="nsew")
        self.yn_var = tk.BooleanVar(value=False)
        yes_btn = tk.Radiobutton(self.lyrics_frame, text="Yes", anchor="w", variable=self.yn_var, value=True, command=self.english_select)
        yes_btn.grid(row=input_row_n, column=1, padx=0, pady=5, sticky="nsew")
        no_btn = tk.Radiobutton(self.lyrics_frame, text="No", anchor="w", variable=self.yn_var, value=False, command=self.english_select)
        no_btn.grid(row=input_row_n, column=2, padx=0, pady=5, sticky="nsew")

        # ----- English Lyrics Frame -----
        self.english_lyrics_frame = tk.Frame(self)
        self.english_lyrics_frame.grid(row=0, column=3, columnspan=3, padx=10, pady=10, sticky="nsew")

        # English lyrics text box
        tk.Label(self.english_lyrics_frame, text="영어 가사 입력").grid(row=0, column=0, columnspan=3)
        eng_lyrics = tk.Text(self.english_lyrics_frame, width=50, height=20)
        eng_lyrics.grid(row=1, column=0, columnspan=3, pady=5)
        self.entry_widgets["english lyrics"] = eng_lyrics

        # English lyrics font size
        tk.Label(self.english_lyrics_frame, text="English Font Size:", anchor="e").grid(row=2, column=0, pady=5, sticky="nsew")
        eng_font_entry = tk.Entry(self.english_lyrics_frame)
        eng_font_entry.insert(0, str(27))
        eng_font_entry.grid(row=2, column=1, columnspan=2, pady=5)
        self.entry_widgets["english font size"] = eng_font_entry

        self.english_lyrics_frame.grid_forget()
        

        # Create presentation button
        create_btn = tk.Button(self, text="파워포인트 생성", command=self.create_presentation)
        create_btn.grid(row=3, columnspan=6, pady=15)

    def english_select(self):
        if self.yn_var.get():
            self.english_lyrics_frame.grid(row=0, column=3, columnspan=3, padx=10, pady=10, sticky="nsew")
        else:
            self.english_lyrics_frame.grid_forget()
    
    def create_presentation(self):
        # get text input
        lyrics_text = self.entry_widgets["lyrics"].get("1.0", tk.END).strip()
        eng_text = self.entry_widgets["english lyrics"].get("1.0", tk.END).strip()
        num_of_lines = int(self.entry_widgets["num of lines"].get().strip())

        if not lyrics_text:
            messagebox.showwarning("Input Error", "Please enter text for the slide.")
            return
        lines = self.parse_lyrics(lyrics_text, num_of_lines)

        eng_lines = []
        if self.yn_var.get() and eng_text:
            eng_lines = self.parse_lyrics(eng_text, num_of_lines)
        
        # get font size
        while True:
            try:
                font_size = int(self.entry_widgets["font size"].get().strip())
                break
            except ValueError:
                messagebox.showwarning("Input Error", "Font size has to be an integer.")
                return
        
        while True:
            try:
                eng_font_size = int(self.entry_widgets["english font size"].get().strip())
                break
            except ValueError:
                messagebox.showwarning("Input Error", "English font size has to be an integer.")
                return
        
        # get font family
        font_family = self.entry_widgets["font family"].get().strip()
        lyrics_position = self.position_var.get()
        print(lyrics_position)

        for idx in range(len(lines)):
        # while(idx < len(lyrics_text)):
            text = lines[idx]
            eng_text = eng_lines[idx] if len(eng_lines) > idx else None
            self.create_slide(text, font_size, font_family, lyrics_position=lyrics_position,eng_lyrics=eng_text, eng_font_size=eng_font_size)
        
        save_path = filedialog.asksaveasfilename(
            defaultextension=".pptx",
            filetypes=[("PowerPoint files", "*.pptx")],
            title="Save Presentation As"
        )

        if not save_path:
            return
            
        # ppt.save('./'+file_name+'.pptx')
        self.ppt.save(save_path)
        self.destroy()

    def parse_lyrics(self, lyrics_text, num_of_lines):
        lines = lyrics_text.splitlines()
        parsed_lyrics = []
        temp = []
        for line in lines:
            if line.strip() == "":
                continue
            temp.append(line)
            if len(temp) == num_of_lines:
                parsed_lyrics.append("\n".join(temp))
                temp = []
        if temp:
            parsed_lyrics.append("\n".join(temp))
        return parsed_lyrics

    def create_slide(self, text, font_size, font_family, lyrics_position, eng_lyrics=None, eng_font_size=None):
        slide_layout = self.ppt.slide_layouts[6]
        slide = self.ppt.slides.add_slide(slide_layout)

        # slide dimensions
        slide_width = self.ppt.slide_width.inches
        slide_height = self.ppt.slide_height.inches

        # text box dimensions
        text_width = Inches(slide_width)
        text_height = Cm(6)
        
        # slide background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0,0,0)

        # centering text box
        left = (slide_width - text_width.inches) / 2

        textBox = slide.shapes.add_textbox(Inches(left), Inches(0), text_width, text_height)
        text_frame = textBox.text_frame

        margin = 0.2  # inches
        if lyrics_position == "Top":
            top = 0
            text_frame.vertical_anchor = MSO_ANCHOR.TOP
        elif lyrics_position == "Bottom":
            top = slide_height - text_height.inches - margin
            text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
        else:  # Center
            top = (slide_height - text_height.inches) / 2
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        textBox.top = Inches(top)

        # Style the text
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = text
        run.font.name = font_family
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(255, 255, 255)

        if eng_lyrics:
            eng_text_height = Cm(6)

            eng_left = 0
            eng_top = slide_height - eng_text_height.inches - margin

            eng_textbox = slide.shapes.add_textbox(
                Inches(eng_left),
                Inches(eng_top),
                Inches(slide_width),
                eng_text_height
            )

            eng_frame = eng_textbox.text_frame
            eng_frame.vertical_anchor = MSO_ANCHOR.BOTTOM

            eng_paragraph = eng_frame.paragraphs[0]
            eng_paragraph.alignment = PP_ALIGN.CENTER

            eng_run = eng_paragraph.add_run()
            eng_run.text = eng_lyrics
            eng_run.font.name = font_family
            eng_run.font.size = Pt(eng_font_size)
            eng_run.font.color.rgb = RGBColor(255, 255, 255)

    def center_window(self):
        self.update_idletasks()  # Ensure widgets are rendered

        window_width = self.winfo_width()
        window_height = self.winfo_height()

        screen_width = self.winfo_screenwidth()
        screen_height = self.winfo_screenheight()

        center_x = (screen_width - window_width) // 2
        center_y = (screen_height - window_height) // 2

        self.geometry(f"+{center_x}+{center_y}")
    
if __name__ == "__main__":
    app = Lyrics()
    app.mainloop()
